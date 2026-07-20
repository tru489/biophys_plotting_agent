"""
biophys_plot_toolkit — reusable plotting library for compiled FXM/SMR/Coulter experiments.

Ported from a hand-written reference analysis and generalized so a short per-experiment driver can
load the raw biophys_helpers outputs and produce the standard figure grid without re-deriving the
plotting internals.

Three layers, low -> high:
  * Loaders           -> records: {sample, props:{name: ndarray}, meta:{col: value, ..., rep}}
                         `meta` carries EVERY annotation column, so any column can drive plots.
  * infer_roles       -> classifies each metadata column into a role (boolean / categorical /
                         time / ordered / continuous / label / structural) so the driver (and
                         Claude) can decide how to group, compare, order and color the data.
  * draw_*            -> low-level primitives: draw on a passed-in axis, no semantics.
  * plot_grouped / compare_groups / timecourse_by / scatter_by / facet / cross_groups
                      -> mid-level combinators parameterized by WHICH column(s) to use.
  * build_plan / render_plan / autoplot
                      -> high-level: infer a plot plan from the roles, show it, execute it.

The loaders read the RAW biophys_helpers outputs directly (no reorg step):
  Coulter — annotate_coulter_samples.py's '*_coulter_sample_annotation/' dir (metadata.csv +
            a single-cell data CSV; columns are samples, rows are per-cell volumes).
            One property "volume" (fL, gated upstream).
  iFXM    — compile_experiment.py's '*_compiled/experiment_data.xlsx' (a 'metadata' sheet + one
            worksheet per sample). A sample's PAIRED ('pair_') block gives, row-aligned:
            mass (pair_mass_pg, pg), density (pair_buoyant_density + baseline_density),
            vol_uncal (pair_volume_au), vol_cal (pair_volume_fL, only when calibrated). A sample
            with no paired block (mass-only / volume-only run) falls back to the standalone MASS
            ('mass_') and/or VOLUME ('vol_') blocks for the distribution props (density needs
            pairing, so it is empty there). scatter (load_ifxm_paired) uses paired samples only.

No statistical outlier rejection is applied anywhere — the loaders drop only non-finite (NaN/inf)
values. The only intentional data exclusions are the metadata-driven gates (bm_gate / ifxm_gate)
and samples skipped because they lack a paired ('pair_') block.

NOTE on baseline_density: the absolute density baseline is variable between experiments and is
NOT stored in any data file, so `load_ifxm` REQUIRES it as an explicit argument (no default).
"""
import re
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
from matplotlib.patches import Patch
from matplotlib.transforms import blended_transform_factory
from pathlib import Path

# ---------------------------------------------------------------------------
# Styling defaults (override in the driver if an experiment needs different keys)
# ---------------------------------------------------------------------------
COND_COLORS = {
    "activated":     "#0072B2",
    "starved":       "#E69F00",
    "drug_treated":  "#009E73",
    "proliferating": "#CC79A7",
}
DRUG_COLORS = {  # keys match the lowercased output of _norm_drug
    "dmso":       "#0072B2",
    "1um-wnk463": "#E69F00",
    "2um-zt-1a":  "#009E73",
}
# High-contrast, colorblind-accommodating pair for booleans: [falsey, truthy] (Okabe-Ito).
BOOL_COLORS = ["#0072B2", "#D55E00"]  # blue (false) vs vermillion (true)
FALLBACK_COLOR = "#999999"

# Colorblind-safe cycle (Okabe-Ito + a few extensions) for auto-assigning colors to unknown values.
_AUTO_PALETTE = [
    "#0072B2", "#E69F00", "#009E73", "#CC79A7", "#D55E00", "#56B4E9",
    "#F0E442", "#000000", "#8C613C", "#666666", "#B8C9D0", "#5975A4",
]

# (prop_key, axis_label) lists — defaults matching the reference experiments.
COULTER_PROPS = [("volume", "Volume (fL)")]
IFXM_PROPS = [
    ("mass",      "Buoyant mass (pg)"),
    ("density",   "Density (g/mL)"),
    ("vol_cal",   "Calibrated volume (fL)"),
    ("vol_uncal", "Volume (fL)"),
]

# Sentinel so a forgotten baseline_density fails loudly rather than silently defaulting.
_REQUIRED = object()


# ---------------------------------------------------------------------------
# Value normalization
# ---------------------------------------------------------------------------
def _is_blank(x) -> bool:
    return x is None or (isinstance(x, float) and np.isnan(x)) or str(x).strip() == ""


def _norm_cond(c) -> str:
    c = "" if _is_blank(c) else str(c).strip()
    return "drug_treated" if c in ("drug_treat", "drug_treated") else c


def _norm_drug(d) -> str:
    if _is_blank(d):
        return ""
    d = str(d).strip().lower()
    return d.replace("zt1a", "zt-1a")


def _rep(sample_name: str) -> str:
    m = re.search(r"rep(\d+)", str(sample_name))
    return f"rep{m.group(1)}" if m else "rep1"


# Applied to matching metadata columns at load time (by column name). Overridable per loader.
# Keeps the reference condition/drug fixups while everything else passes through stripped.
VALUE_NORMALIZERS = {"condition": _norm_cond, "drug_name": _norm_drug}


def _clean_value(v):
    """Light normalization for a raw metadata cell: blanks -> '', strings stripped, else as-is."""
    if _is_blank(v):
        return ""
    return v.strip() if isinstance(v, str) else v


def _build_meta(row, sample, normalizers) -> dict:
    """The generic annotation bag for one record: every column, lightly normalized, plus rep."""
    meta = {}
    for col in row.index:
        val = _clean_value(row[col])
        fn = normalizers.get(col)
        meta[col] = fn(val) if fn else val
    meta["rep"] = _rep(sample)
    return meta


# ---------------------------------------------------------------------------
# Robust metadata access
# ---------------------------------------------------------------------------
def rget(r, col, default=None):
    """Value of annotation `col` for record `r` (from its meta bag), or `default`."""
    v = r.get("meta", {}).get(col, default)
    return default if (col != "rep" and _is_blank(v)) else v


def _get(row, col, default=""):
    if col and col in row.index and not _is_blank(row[col]):
        return row[col]
    return default


def _get_float(row, col, default):
    v = _get(row, col, None)
    try:
        return float(v) if v is not None else default
    except (TypeError, ValueError):
        return default


def _gate_bound(row, col, default):
    """A gate bound as float, or `default` (±inf) when the column is missing or NaN, so a
    missing/ungated experiment means 'no cutoff' rather than dropping every cell."""
    if col and col in row.index and not _is_blank(row[col]):
        try:
            return float(row[col])
        except (TypeError, ValueError):
            return default
    return default


def _require_col(meta, col, role):
    if col not in meta.columns:
        raise KeyError(
            f"could not find the {role} column '{col}'. Available columns: "
            f"{list(meta.columns)}. Pass the right name via the loader's *_col argument."
        )


def _to_floats(values):
    """Return values as a list of floats if every non-blank value is numeric, else None."""
    out = []
    for v in values:
        if _is_blank(v):
            continue
        try:
            out.append(float(v))
        except (TypeError, ValueError):
            return None
    return out


# ---------------------------------------------------------------------------
# Color maps
# ---------------------------------------------------------------------------
def build_color_map(values, base=None):
    """Color map covering every value in `values`. Known keys in `base` keep their color; unknown
    values get the next distinct palette color (deterministic, order-stable)."""
    out = dict(base or {})
    used = set(out.values())
    i = 0
    for v in values:
        if v in out:
            continue
        while i < len(_AUTO_PALETTE) and _AUTO_PALETTE[i] in used:
            i += 1
        color = _AUTO_PALETTE[i % len(_AUTO_PALETTE)]
        out[v] = color
        used.add(color)
        i += 1
    return out


def color_map_for(records, col, roles=None, base=None):
    """Color map for the values of `col` present in `records`. Booleans get the high-contrast
    BOOL_COLORS pair; `condition`/`drug`-named columns seed from COND_COLORS / DRUG_COLORS."""
    role = _role_for(records, col, roles)
    values = group_order(records, col, roles)
    if role["role"] == "boolean" and len(values) <= 2:
        # falsey -> BOOL_COLORS[0], truthy -> BOOL_COLORS[1] (values already ordered falsey->truthy)
        return {v: BOOL_COLORS[i] for i, v in enumerate(values)}
    if base is None:
        if col == "condition":
            base = COND_COLORS
        elif "drug" in col.lower():
            base = DRUG_COLORS
    return build_color_map(values, base)


# ---------------------------------------------------------------------------
# Loaders  ->  list of records {sample, props:{name: arr}, meta:{col: value, ..., rep}}
# ---------------------------------------------------------------------------

# Columns inside each iFXM sample sheet's blocks, after the block prefix is stripped. This is the
# compile_experiment.py output contract. Each sample sheet holds up to three side-by-side blocks:
#   PAIRED ('pair_')  — matched cells, row-aligned: mass_pg, buoyant_density, volume_au, volume_fL.
#   MASS   ('mass_')  — every SMR cell (unpaired): mass_pg (+ pass-through mass_* columns).
#   VOLUME ('vol_')   — every FXM cell (unpaired): volume_au, volume_fL.
# Density is pairing-only; the standalone blocks never carry it.
_PAIR_MASS = "mass_pg"          # buoyant mass (pg)          (was 'matched_mass' in the old h5)
_PAIR_DENS = "buoyant_density"  # RELATIVE density (g/mL); absolute = + baseline_density
_PAIR_VUN  = "volume_au"        # uncalibrated volume (AU)   (was 'volume' in the old h5)
_PAIR_VCAL = "volume_fL"        # calibrated volume (fL); present ONLY when a calibration ran
# Standalone-block value columns (after the 'mass_'/'vol_' prefix is stripped):
_MASS_STANDALONE = "mass_pg"    # MASS block buoyant mass (pg)   -> 'mass_mass_pg' on the sheet
_VOL_UNCAL = "volume_au"        # VOLUME block uncalibrated (AU) -> 'vol_volume_au'
_VOL_CAL   = "volume_fL"        # VOLUME block calibrated (fL)   -> 'vol_volume_fL' (if calibrated)


def _find_coulter_data_csv(d: Path, meta_path: Path) -> Path:
    """The single-cell data CSV keeps the input file's name, so it is the one CSV in the
    directory that is not metadata.csv. Raise clearly if it is ambiguous or missing."""
    csvs = [p for p in sorted(d.glob("*.csv")) if p.resolve() != meta_path.resolve()]
    if len(csvs) == 1:
        return csvs[0]
    if not csvs:
        raise FileNotFoundError(
            f"no single-cell data CSV found in {d} (only metadata.csv). Pass data_file=.")
    raise ValueError(
        f"multiple candidate data CSVs in {d}: {[p.name for p in csvs]}; pass data_file= to pick one.")


def load_coulter(coulter_dir, *, sample_col="sample_name", data_file=None,
                 normalizers=VALUE_NORMALIZERS) -> list:
    """Load Coulter single-cell volumes from a '*_coulter_sample_annotation/' directory produced
    by annotate_coulter_samples.py (or point straight at its metadata.csv).

    metadata.csv has `sample_name` plus whatever annotation columns were hand-added. The per-cell
    data is a separate CSV whose COLUMNS are samples (headers == sample_name values) and whose ROWS
    are single-cell volumes; it is auto-located as the non-metadata CSV in the dir (`data_file`
    overrides). Volume units follow the input CSV (fL in the standard Coulter pipeline).

    Every metadata column is carried into each record's `meta` bag (role assignment happens later
    via infer_roles), so only `sample_col` is required.
    """
    p = Path(coulter_dir)
    if p.is_file():
        meta_path, d = p, p.parent
    else:
        meta_path, d = p / "metadata.csv", p
    if not meta_path.exists():
        raise FileNotFoundError(
            f"metadata.csv not found at {meta_path}. Point load_coulter at the "
            f"'*_coulter_sample_annotation/' directory from annotate_coulter_samples.py.")
    meta = pd.read_csv(meta_path)
    _require_col(meta, sample_col, "sample-name")

    data_path = Path(data_file) if data_file else _find_coulter_data_csv(d, meta_path)
    data = pd.read_csv(data_path)

    recs = []
    for _, r in meta.iterrows():
        sample = _get(r, sample_col, "")
        if sample not in data.columns:
            raise KeyError(
                f"sample {sample!r} has no column in {data_path.name}. "
                f"Data columns: {list(data.columns)[:12]}{' ...' if data.shape[1] > 12 else ''}")
        arr = data[sample].to_numpy(dtype=float)
        arr = arr[np.isfinite(arr)]  # drop NaN padding / non-finite; no statistical outlier rejection
        recs.append({
            "sample": sample,
            "props":  {"volume": arr},
            "meta":   _build_meta(r, sample, normalizers),
        })
    return recs


def _open_ifxm_xlsx(compiled_dir) -> Path:
    """Resolve experiment_data.xlsx from a '*_compiled/' dir (or accept the .xlsx path directly)."""
    p = Path(compiled_dir)
    xlsx = p if p.suffix.lower() == ".xlsx" else p / "experiment_data.xlsx"
    if not xlsx.exists():
        raise FileNotFoundError(
            f"experiment_data.xlsx not found at {xlsx}. Point load_ifxm at the "
            f"'*_compiled/' directory produced by compile_experiment.py.")
    return xlsx


def _read_block(xls, xlsx_path: Path, sheet_name, prefix, _sheet_cache=None) -> pd.DataFrame:
    """One prefixed block ('pair'|'mass'|'vol') of a sample sheet as a DataFrame (rows dropped where
    all block columns are NaN, prefix stripped), or None if that block is absent. A block that
    overflowed Excel's row limit is written by compile_experiment.py to a sibling
    '{sheet}_{prefix}_overflow.csv'; that full copy is preferred when present. `_sheet_cache` is an
    optional {sheet_name: DataFrame} dict so the three blocks of a sheet share one read_excel."""
    overflow = Path(xlsx_path).parent / f"{sheet_name}_{prefix}_overflow.csv"
    if overflow.exists():
        blk = pd.read_csv(overflow)
    else:
        if _sheet_cache is not None and sheet_name in _sheet_cache:
            sheet = _sheet_cache[sheet_name]
        else:
            sheet = pd.read_excel(xls, sheet_name=sheet_name)
            if _sheet_cache is not None:
                _sheet_cache[sheet_name] = sheet
        blk = sheet.filter(regex=rf"^{prefix}_").dropna(how="all")
    if blk.empty:
        return None
    pre = f"{prefix}_"
    blk = blk.rename(columns=lambda c: c[len(pre):] if str(c).startswith(pre) else c)
    return blk


_EMPTY = np.array([], dtype=float)


def _require_baseline(baseline_density):
    if baseline_density is _REQUIRED:
        raise ValueError(
            "baseline_density (g/mL) is required to compute density from a paired block — it varies "
            "per experiment and is not stored in the data files. Set it in your driver. (It can be "
            "omitted only for a mass-only / volume-only experiment with no paired iFXM data.)")
    return baseline_density


def _gate_clean(a, mask):
    """Apply a keep-mask (if lengths match) then drop non-finite. Empty arrays pass through."""
    a = a[mask] if a.size == mask.size else a
    return a[np.isfinite(a)]


def _load_ifxm_records(compiled_dir, baseline_density, sample_col, sheet_col, gate_cols,
                       normalizers, paired):
    """Shared iFXM reader. `paired`=True keeps all props row-aligned under one mask from the PAIRED
    block (for scatter); unpaired samples are skipped. `paired`=False builds distribution props:
    from the PAIRED block when present (matched subset — unchanged), else falling back to the
    standalone MASS / VOLUME blocks so mass-only and volume-only runs still load (density stays
    pairing-only)."""
    bm_lo_c, bm_hi_c, ix_lo_c, ix_hi_c = gate_cols
    xlsx = _open_ifxm_xlsx(compiled_dir)
    recs = []
    with pd.ExcelFile(xlsx) as xls:
        meta = pd.read_excel(xls, sheet_name="metadata")
        _require_col(meta, sample_col, "sample-name")
        skey = sheet_col if sheet_col in meta.columns else sample_col
        cache = {}

        for _, r in meta.iterrows():
            pair = _read_block(xls, xlsx, r[skey], "pair", cache)
            has_pair = pair is not None and _PAIR_MASS in pair.columns and _PAIR_VUN in pair.columns

            if paired and not has_pair:
                continue  # scatter needs a paired block; unpaired samples have no row-aligned pairs

            bm_lo = _gate_bound(r, bm_lo_c, -np.inf)
            bm_hi = _gate_bound(r, bm_hi_c,  np.inf)
            ix_lo = _gate_bound(r, ix_lo_c, -np.inf)
            ix_hi = _gate_bound(r, ix_hi_c,  np.inf)
            sample = _get(r, sample_col, "")
            meta_bag = _build_meta(r, sample, normalizers)

            if has_pair:
                mass = pair[_PAIR_MASS].to_numpy(dtype=float)
                dens = pair[_PAIR_DENS].to_numpy(dtype=float) + _require_baseline(baseline_density)
                vun  = pair[_PAIR_VUN].to_numpy(dtype=float)
                has_cal = _PAIR_VCAL in pair.columns
                vcal = pair[_PAIR_VCAL].to_numpy(dtype=float) if has_cal else _EMPTY

                if paired:
                    # single mask over always-present props keeps them length-matched & row-aligned;
                    # vcal (finite wherever volume_au is) rides along under the same mask.
                    mask = (np.isfinite(mass) & np.isfinite(dens) & np.isfinite(vun)
                            & (mass >= bm_lo) & (mass <= bm_hi)
                            & (vun >= ix_lo) & (vun <= ix_hi))
                    props = {
                        "mass":      mass[mask],
                        "density":   dens[mask],
                        "vol_cal":   vcal[mask] if has_cal else _EMPTY,
                        "vol_uncal": vun[mask],
                    }
                else:
                    bm_mask = np.isfinite(mass) & (mass >= bm_lo) & (mass <= bm_hi)
                    ix_mask = np.isfinite(vun) & (vun >= ix_lo) & (vun <= ix_hi)
                    props = {
                        "mass":      _gate_clean(mass, bm_mask),
                        "density":   _gate_clean(dens, ix_mask),
                        "vol_cal":   _gate_clean(vcal, ix_mask),
                        "vol_uncal": _gate_clean(vun, ix_mask),
                    }
            else:
                # Unpaired sample (mass-only / volume-only): fall back to the standalone blocks.
                massblk = _read_block(xls, xlsx, r[skey], "mass", cache)
                volblk  = _read_block(xls, xlsx, r[skey], "vol", cache)
                if massblk is None and volblk is None:
                    continue  # sample has no tabular iFXM data at all

                mass = massblk[_MASS_STANDALONE].to_numpy(dtype=float) \
                    if massblk is not None and _MASS_STANDALONE in massblk.columns else _EMPTY
                vun = volblk[_VOL_UNCAL].to_numpy(dtype=float) \
                    if volblk is not None and _VOL_UNCAL in volblk.columns else _EMPTY
                vcal = volblk[_VOL_CAL].to_numpy(dtype=float) \
                    if volblk is not None and _VOL_CAL in volblk.columns else _EMPTY

                bm_mask = np.isfinite(mass) & (mass >= bm_lo) & (mass <= bm_hi)
                ix_mask = np.isfinite(vun) & (vun >= ix_lo) & (vun <= ix_hi)
                props = {
                    "mass":      _gate_clean(mass, bm_mask),
                    "density":   _EMPTY,               # density requires pairing
                    "vol_cal":   _gate_clean(vcal, ix_mask),
                    "vol_uncal": _gate_clean(vun, ix_mask),
                }
            recs.append({"sample": sample, "props": props, "meta": meta_bag})
    return recs


def load_ifxm(compiled_dir, baseline_density=_REQUIRED, *, sample_col="sample_name",
              sheet_col="sheet_name", bm_lower_col="bm_gate_lower", bm_upper_col="bm_gate_upper",
              ifxm_lower_col="ifxm_gate_lower", ifxm_upper_col="ifxm_gate_upper",
              normalizers=VALUE_NORMALIZERS) -> list:
    """Load iFXM distribution data from a '*_compiled/' dir's experiment_data.xlsx.

    baseline_density: fluid baseline added to the measured RELATIVE buoyant density to get absolute
    density (g/mL). Not stored in the data — supply it whenever the experiment has paired iFXM data.
    It is required lazily: only raises if a PAIRED block is actually read, so a mass-only /
    volume-only experiment can omit it.

    For each sample (worksheet named by `sheet_col`): if it has a PAIRED ('pair_') block, mass /
    density / vol_cal / vol_uncal come from that matched subset (mass bm-gated; the others ifxm-gated
    on the uncalibrated volume). If it has NO paired block (a mass-only or volume-only run), the
    standalone MASS ('mass_') and/or VOLUME ('vol_') blocks are used instead — `mass` from the full
    SMR distribution and/or `vol_uncal`/`vol_cal` from the full FXM distribution, with `density`
    empty (density requires pairing). Samples with no tabular iFXM data are skipped. Every metadata
    column is carried into each record's `meta`.
    """
    return _load_ifxm_records(
        compiled_dir, baseline_density, sample_col, sheet_col,
        (bm_lower_col, bm_upper_col, ifxm_lower_col, ifxm_upper_col), normalizers, paired=False)


def load_ifxm_paired(compiled_dir, baseline_density=_REQUIRED, *, sample_col="sample_name",
                     sheet_col="sheet_name", bm_lower_col="bm_gate_lower",
                     bm_upper_col="bm_gate_upper", ifxm_lower_col="ifxm_gate_lower",
                     ifxm_upper_col="ifxm_gate_upper", normalizers=VALUE_NORMALIZERS) -> list:
    """Like load_ifxm, but keeps per-cell arrays row-ALIGNED across properties (one shared mask from
    the PAIRED block), so a cell's mass / density / volume stay paired. Use for scatter_by. Only
    samples with a paired block appear (unpaired mass-only / volume-only samples are skipped — there
    is nothing to correlate); samples not calibrated get an empty `vol_cal` (scatters using it are
    skipped, not misaligned). baseline_density is always needed here (density is always read)."""
    return _load_ifxm_records(
        compiled_dir, baseline_density, sample_col, sheet_col,
        (bm_lower_col, bm_upper_col, ifxm_lower_col, ifxm_upper_col), normalizers, paired=True)


# ---------------------------------------------------------------------------
# Role inference — classify each metadata column so plots can be chosen intelligently
# ---------------------------------------------------------------------------
DEFAULT_STRUCTURAL = {"sheet_name", "hdf5_key", "coulter_column", "calibration_factor"}
STRUCTURAL_PATTERNS = [re.compile(p) for p in (r"^has_", r"_gate_lower$", r"_gate_upper$",
                                               r"^bm_gate", r"^ifxm_gate")]
IDENTITY_COLS = {"sample_name"}

_TIME_NAME = re.compile(r"(?i)(^|_)(t|time|elapsed)(_|$)")
# suffix -> unit, longest first so '_hours' wins over '_h'
_TIME_SUFFIX = [("_hours", "h"), ("_hour", "h"), ("_hrs", "h"), ("_hr", "h"), ("_h", "h"),
                ("_minutes", "min"), ("_mins", "min"), ("_min", "min"),
                ("_seconds", "s"), ("_sec", "s")]
_UNIT_PER_HOUR = {"h": 1.0, "min": 60.0, "s": 3600.0}
_GRADIENT_HINTS = re.compile(
    r"(?i)(dose|conc|concentration|passage|day|cycle|generation|gen|temp|ph|dilution)")

_BOOL_SETS = [{"yes", "no"}, {"true", "false"}, {"0", "1"}, {"y", "n"}, {"t", "f"}]
_TRUTHY = {"yes", "true", "1", "y", "t"}
_ORDERED_MAX_CARD = 15   # numeric non-time non-hint below this is proposed as an ordered gradient


def _is_structural(col: str) -> bool:
    return col in DEFAULT_STRUCTURAL or any(p.search(col) for p in STRUCTURAL_PATTERNS)


def _time_unit(col: str):
    """(unit, confirm) for a time column name, or (None, None) if the name is not time-like.
    confirm=True when the unit had to be assumed (generic 'time'/'t' with no unit suffix)."""
    c = col.lower()
    for suf, u in _TIME_SUFFIX:
        if c.endswith(suf):
            return u, False
    if _TIME_NAME.search(col):
        return "h", True   # generic time name, assume hours but flag for confirmation
    return None, None


def _role(col, role, values, order, **kw):
    info = {"col": col, "role": role, "values": list(values), "order": list(order),
            "unit": None, "to_hours": None, "cardinality": len(set(values)),
            "use_for": set(), "reason": "", "confirm": False}
    info.update(kw)
    return info


def classify_column(col, values) -> dict:
    """Classify a single column from its NAME and its non-blank VALUES into a RoleInfo dict.
    Pure and dtype-agnostic (numeric-ness is inferred by float-coercion), so it works from either a
    metadata DataFrame column or the values in loaded records."""
    vals = [v for v in values if not _is_blank(v)]
    uniq = list(dict.fromkeys(vals))                 # unique, order-preserving
    low = {str(v).strip().lower() for v in uniq}
    nums = _to_floats(uniq)

    if _is_structural(col):
        return _role(col, "structural", uniq, uniq, reason="pipeline/structural column — ignored")
    if col in IDENTITY_COLS:
        return _role(col, "label", uniq, uniq, use_for={"label"}, reason="sample identity")

    # boolean (checkbox columns are literally 'yes'/'no')
    if low and any(low <= s for s in _BOOL_SETS) and len(low) <= 2:
        order = sorted(uniq, key=lambda v: str(v).strip().lower() in _TRUTHY)  # falsey -> truthy
        return _role(col, "boolean", uniq, order,
                     use_for={"group", "compare", "facet", "color", "series"},
                     reason=f"boolean ({'/'.join(map(str, order))})")

    # time
    unit, confirm = _time_unit(col)
    if nums is not None and unit is not None:
        per_h = _UNIT_PER_HOUR[unit]
        order = sorted(uniq, key=lambda v: float(v))
        return _role(col, "time", uniq, order, unit=unit, to_hours=(lambda v, k=per_h: float(v) / k),
                     use_for={"order", "series", "color"}, confirm=confirm,
                     reason=f"time in {unit}" + (" (unit assumed — confirm)" if confirm else ""))

    # ordered / gradient (numeric, non-time)
    if nums is not None:
        hinted = bool(_GRADIENT_HINTS.search(col))
        if hinted or len(uniq) <= _ORDERED_MAX_CARD:
            order = sorted(uniq, key=lambda v: float(v))
            why = "named like a gradient" if hinted else f"{len(uniq)} distinct numeric values"
            return _role(col, "ordered", uniq, order,
                         use_for={"order", "group", "compare", "color", "series"}, confirm=True,
                         reason=f"numeric gradient? ({why}) — confirm ordered vs categorical")
        return _role(col, "continuous", uniq, sorted(uniq, key=lambda v: float(v)),
                     use_for={"color"}, reason="continuous numeric — color/scatter axis only")

    # categorical (any cardinality; no cap)
    order = sorted(uniq, key=lambda v: (-vals.count(v), str(v)))   # frequency desc, then alpha
    note = "" if len(uniq) < len(vals) else " (one value per sample)"
    return _role(col, "categorical", uniq, order,
                 use_for={"group", "compare", "facet", "color", "series"},
                 reason=f"categorical, {len(uniq)} values{note}")


def infer_roles(records, *, overrides=None, skip=("rep",)) -> dict:
    """Classify every metadata column present across `records` into a RoleInfo. `overrides` is a
    {col: role_name} map that pins a column's role (e.g. force a numeric column to 'ordered' or
    'categorical' after the user confirms the plan). `skip` omits synthesized columns from the
    report (rep is always available for ordering/labels regardless)."""
    cols = []
    for r in records:
        for c in r.get("meta", {}):
            if c not in cols:
                cols.append(c)
    roles = {}
    for col in cols:
        if col in skip:
            continue
        info = classify_column(col, [rget(r, col) for r in records])
        ov = (overrides or {}).get(col)
        if ov and ov != info["role"]:
            info = classify_column(col, [rget(r, col) for r in records])  # recompute base
            info["role"] = ov
            info["confirm"] = False
            info["reason"] = f"role overridden to {ov}"
            info["use_for"] = {
                "boolean": {"group", "compare", "facet", "color", "series"},
                "categorical": {"group", "compare", "facet", "color", "series"},
                "ordered": {"order", "group", "compare", "color", "series"},
                "time": {"order", "series", "color"},
                "continuous": {"color"},
                "label": {"label"}, "structural": set(),
            }.get(ov, info["use_for"])
        roles[col] = info
    return roles


def _role_for(records, col, roles):
    """RoleInfo for `col`: from `roles` if given, else inferred on the fly from record values."""
    if roles and col in roles:
        return roles[col]
    return classify_column(col, [rget(r, col) for r in records])


# ---------------------------------------------------------------------------
# Ordering + labels
# ---------------------------------------------------------------------------
def group_order(records, col, roles=None) -> list:
    """Values of `col` present in `records`, in canonical plotting order (role-defined)."""
    role = _role_for(records, col, roles)
    present = list(dict.fromkeys(rget(r, col) for r in records if rget(r, col) is not None))
    ordered = [v for v in role["order"] if v in present]
    ordered += [v for v in present if v not in ordered]
    return ordered


def _sort_key(r, col, role):
    v = rget(r, col)
    if role["role"] in ("time", "ordered", "continuous"):
        try:
            return (0, float(v))
        except (TypeError, ValueError):
            return (1, 0.0)
    order = role["order"]
    return (order.index(v), 0.0) if v in order else (len(order), 0.0)


def sort_records(records, by, roles=None) -> list:
    """Sort records by a list of (col, ...) — each col ordered per its role (time/ordered numeric,
    else canonical categorical order). `by` may be a list of column names or (col, _) pairs."""
    cols = [b[0] if isinstance(b, (tuple, list)) else b for b in by]
    role_map = {c: _role_for(records, c, roles) for c in cols}

    def key(r):
        return tuple(_sort_key(r, c, role_map[c]) for c in cols)
    return sorted(records, key=key)


def _time_label(t) -> str:
    try:
        t = float(t)
    except (TypeError, ValueError):
        return str(t)
    return f"{int(t)}h" if t == int(t) else f"{t}h"


def value_label(col, value, roles=None, records=None) -> str:
    """Short label for a single value of `col`, role-aware (time -> '6h', boolean -> 'is_x=yes')."""
    role = roles.get(col) if roles and col in roles else (
        _role_for(records, col, None) if records is not None else {"role": "categorical", "unit": None})
    if role["role"] == "time":
        conv = role.get("to_hours")
        return _time_label(conv(value) if conv else value)
    if role["role"] == "boolean":
        return f"{col}={value}"
    return str(value)


def _detail_label(r, roles, time_col) -> str:
    """Per-sample row/box label inside a single group: time (if any) + rep, else the sample name."""
    parts = []
    if time_col is not None and rget(r, time_col) is not None:
        parts.append(value_label(time_col, rget(r, time_col), roles))
    parts.append(rget(r, "rep"))
    return " ".join(p for p in parts if p) or r["sample"]


def _compare_label(r, group_col, roles, time_col) -> str:
    """Per-sample label in a cross-group comparison: group value + time + rep."""
    parts = [str(rget(r, group_col))]
    if time_col is not None and rget(r, time_col) is not None:
        parts.append(value_label(time_col, rget(r, time_col), roles))
    parts.append(rget(r, "rep"))
    return " ".join(p for p in parts if p)


# ---------------------------------------------------------------------------
# Low-level primitives (operate on a passed-in ax; no semantics)
# ---------------------------------------------------------------------------
def _run_separators(ax, keys, label_fn=None) -> None:
    """Bold labels under the axis with dark vertical separators between runs of equal `keys`."""
    trans = blended_transform_factory(ax.transData, ax.transAxes)
    n = len(keys)
    prev, start = object(), 0
    for i, k in enumerate(list(keys) + [object()]):
        if i == n or k != prev:
            if i > 0 and start < n:
                mid = (start + i - 1) / 2
                lab = label_fn(prev) if label_fn else str(prev)
                ax.text(mid, -0.30, lab, ha="center", va="top", transform=trans,
                        fontsize=9, fontweight="bold")
                if i < n:
                    ax.axvline(i - 0.5, color="black", lw=1.2, alpha=0.8, zorder=1)
            start, prev = i, k


def draw_ridge(ax, arrays, labels, colors, xlabel, overlap: float = 1.7) -> None:
    """Ridge plot of per-row histograms (shared bins, max-normalized), stacked top-down."""
    lo = min(v.min() for v in arrays)
    hi = max(v.max() for v in arrays)
    bins = np.linspace(lo, hi, 41)
    n = len(arrays)
    for i, (vals, c) in enumerate(zip(arrays, colors)):
        counts, _ = np.histogram(vals, bins=bins, density=True)
        h = counts / counts.max() * overlap if counts.max() > 0 else counts
        stair = np.concatenate([h, h[-1:]])
        base = n - 1 - i
        ax.fill_between(bins, base, base + stair, step="post", color=c, alpha=0.6, zorder=n - i)
        ax.step(bins, base + stair, where="post", color="black", lw=0.8, zorder=n - i)
    ax.set_yticks([n - 1 - i for i in range(n)])
    ax.set_yticklabels(labels, fontsize=7)
    ax.set_xlabel(xlabel)


def draw_boxes(ax, arrays, labels, colors, ylabel, sep_keys=None, sep_label_fn=None) -> None:
    """Boxplot + jittered datapoints, one box per array. Optional run separators under the axis."""
    n = len(arrays)
    for i, (vals, c) in enumerate(zip(arrays, colors)):
        jitter = np.random.uniform(-0.18, 0.18, len(vals))
        ax.scatter(np.full(len(vals), i, float) + jitter, vals,
                   color=c, alpha=0.1, s=4, zorder=2, linewidths=0)
        ax.boxplot(vals, positions=[i], widths=0.5, patch_artist=True, showfliers=False,
                   boxprops=dict(facecolor=c, alpha=0.5),
                   medianprops=dict(color="black", linewidth=1.5),
                   whiskerprops=dict(color=c), capprops=dict(color=c))
    ax.set_xticks(range(n))
    ax.set_xticklabels(labels, fontsize=7, rotation=45, ha="right")
    if sep_keys is not None:
        _run_separators(ax, sep_keys, sep_label_fn)
    ax.set_xlim(-0.5, n - 0.5)
    ax.set_ylabel(ylabel)


def draw_ecdf(ax, arrays, labels, colors, xlabel) -> None:
    """Overlaid empirical CDFs, one line per array."""
    for vals, c, lab in zip(arrays, colors, labels):
        x = np.sort(np.asarray(vals, float))
        if x.size == 0:
            continue
        y = np.arange(1, x.size + 1) / x.size
        ax.plot(x, y, color=c, lw=1.6, label=str(lab))
    ax.set_xlabel(xlabel)
    ax.set_ylabel("Cumulative fraction")
    ax.legend(frameon=False, fontsize=8)


def draw_timecourse(ax, series: dict, colors: dict, xlabel, ylabel) -> None:
    """Per-replicate means as points + per-x average as a line, one series per key.
    `series`: {key: [(x, y), ...]}."""
    for sval, pts in sorted(series.items(), key=lambda kv: str(kv[0])):
        c = colors.get(sval, FALLBACK_COLOR)
        t = np.array([p[0] for p in pts], float)
        y = np.array([p[1] for p in pts], float)
        ax.scatter(t, y, color=c, s=30, alpha=0.7, zorder=3)
        uniq = sorted(set(t))
        avg = [y[t == u].mean() for u in uniq]
        ax.plot(uniq, avg, color=c, lw=1.5, marker="o", ms=6, zorder=2, label=str(sval))
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.legend(frameon=False, fontsize=8)


def draw_scatter_marginal(fig, subplot_spec, x, y, color, xlabel, ylabel, title) -> None:
    """A per-cell scatter with marginal histograms (x on top, y on right), inside subplot_spec."""
    inner = gridspec.GridSpecFromSubplotSpec(
        2, 2, subplot_spec=subplot_spec, width_ratios=[3, 1], height_ratios=[1, 3],
        hspace=0.03, wspace=0.03)
    ax_sc = fig.add_subplot(inner[1, 0])
    ax_xh = fig.add_subplot(inner[0, 0], sharex=ax_sc)
    ax_yh = fig.add_subplot(inner[1, 1], sharey=ax_sc)
    fig.add_subplot(inner[0, 1]).axis("off")

    ax_sc.scatter(x, y, color=color, s=2, alpha=0.3, linewidths=0)
    ax_xh.hist(x, bins=30, color=color, alpha=0.7)
    ax_yh.hist(y, bins=30, orientation="horizontal", color=color, alpha=0.7)

    plt.setp(ax_xh.get_xticklabels(), visible=False)
    plt.setp(ax_yh.get_yticklabels(), visible=False)
    ax_xh.tick_params(bottom=False, labelsize=6)
    ax_yh.tick_params(left=False, labelsize=6)
    ax_xh.set_title(title, fontsize=8)
    ax_sc.set_xlabel(xlabel, fontsize=7)
    ax_sc.set_ylabel(ylabel, fontsize=7)
    ax_sc.tick_params(labelsize=6)


# Backward-compatible private aliases (older code referenced the underscore names).
_ridge, _boxes = draw_ridge, draw_boxes


# ---------------------------------------------------------------------------
# Small utilities
# ---------------------------------------------------------------------------
def _slug(value) -> str:
    return re.sub(r"[^0-9a-zA-Z]+", "-", str(value).strip().lower()).strip("-") or "na"


def _save(fig, name: str, out_dir) -> None:
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_dir / name, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  saved {out_dir.name}/{name}")


def _with_prop(records: list, prop: str) -> list:
    return [r for r in records if len(r["props"].get(prop, [])) > 0]


def _find_time_col(records, roles):
    """The first column whose role is 'time', if any (used to order samples within groups)."""
    roles = roles or infer_roles(records)
    for col, info in roles.items():
        if info["role"] == "time":
            return col
    return None


# ---------------------------------------------------------------------------
# Mid-level combinators (parameterized by WHICH column to group / compare / order by)
# ---------------------------------------------------------------------------
def plot_grouped(records, prop, ylabel, datatype, fig_dir, *, group_col, roles=None,
                 kinds=("ridge", "box"), time_col=None, colors=None) -> None:
    """DETAIL within each group: one figure per value of `group_col`; each ridge row / box is a
    sample in that group, ordered by time (if a time column exists) then rep.
    Files: {datatype}_{prop}_{kind}_{group_col}={slug(value)}.png"""
    recs = _with_prop(records, prop)
    if not recs:
        return
    roles = roles or infer_roles(records)
    time_col = time_col if time_col is not None else _find_time_col(records, roles)
    colors = colors or color_map_for(records, group_col, roles)
    for v in group_order(recs, group_col, roles):
        sub = [r for r in recs if rget(r, group_col) == v]
        order_by = [c for c in (time_col, "rep") if c]
        sub = sort_records(sub, order_by, roles) if order_by else sub
        if not sub:
            continue
        arrays = [r["props"][prop] for r in sub]
        col = colors.get(v, FALLBACK_COLOR)
        labels = [_detail_label(r, roles, time_col) for r in sub]
        title = f"{datatype} {prop} — {value_label(group_col, v, roles)}"
        tag = f"{group_col}={_slug(v)}"
        if "ridge" in kinds:
            fig, ax = plt.subplots(figsize=(8, max(3, len(sub) * 0.55)))
            draw_ridge(ax, arrays, labels, [col] * len(sub), ylabel)
            ax.set_title(title)
            _save(fig, f"{datatype}_{prop}_ridge_{tag}.png", fig_dir)
        if "box" in kinds:
            times = [rget(r, time_col) for r in sub] if time_col else None
            fig, ax = plt.subplots(figsize=(max(5, len(sub) * 0.9), 5))
            draw_boxes(ax, arrays, labels, [col] * len(sub), ylabel,
                       sep_keys=times, sep_label_fn=_time_label if times else None)
            ax.set_title(title)
            _save(fig, f"{datatype}_{prop}_box_{tag}.png", fig_dir)


def compare_groups(records, prop, ylabel, datatype, fig_dir, *, group_col, roles=None,
                   kinds=("box", "ridge"), agg="per_sample", colors=None, time_col=None) -> None:
    """COMPARISON across the values of `group_col`. agg='per_sample' (default) draws one box/ridge
    row per sample, colored by its group value and separated by group; agg='pool' pools all cells
    per group value into a single box/ridge row. Both a box and a ridge are produced by default.
    Files: {datatype}_{prop}_{kind}_by_{group_col}.png"""
    recs = _with_prop(records, prop)
    if not recs:
        return
    roles = roles or infer_roles(records)
    time_col = time_col if time_col is not None else _find_time_col(records, roles)
    colors = colors or color_map_for(records, group_col, roles)
    values = group_order(recs, group_col, roles)
    if len(values) < 1:
        return

    if agg == "pool":
        arrays, labels, cols = [], [], []
        for v in values:
            pooled = np.concatenate([r["props"][prop] for r in recs if rget(r, group_col) == v])
            if len(pooled) == 0:
                continue
            arrays.append(pooled)
            labels.append(value_label(group_col, v, roles))
            cols.append(colors.get(v, FALLBACK_COLOR))
        sep_keys = None
        handles = None
    else:  # per_sample
        sub = sort_records(recs, [c for c in (group_col, time_col, "rep") if c], roles)
        arrays = [r["props"][prop] for r in sub]
        labels = [_compare_label(r, group_col, roles, time_col) for r in sub]
        cols = [colors.get(rget(r, group_col), FALLBACK_COLOR) for r in sub]
        sep_keys = [rget(r, group_col) for r in sub]
        handles = [Patch(facecolor=colors.get(v, FALLBACK_COLOR),
                         label=value_label(group_col, v, roles)) for v in values]
    if not arrays:
        return

    title = f"{datatype} {prop} — by {group_col}"
    if "box" in kinds:
        fig, ax = plt.subplots(figsize=(max(5, len(arrays) * 0.9), 5))
        draw_boxes(ax, arrays, labels, cols, ylabel,
                   sep_keys=sep_keys, sep_label_fn=(lambda v: str(v)) if sep_keys else None)
        if handles:
            ax.legend(handles=handles, frameon=False, fontsize=8)
        ax.set_title(title)
        _save(fig, f"{datatype}_{prop}_box_by_{group_col}.png", fig_dir)
    if "ridge" in kinds:
        fig, ax = plt.subplots(figsize=(9, max(3, len(arrays) * 0.55)))
        draw_ridge(ax, arrays, labels, cols, ylabel)
        if handles:
            ax.legend(handles=handles, frameon=False, fontsize=8)
        ax.set_title(title)
        _save(fig, f"{datatype}_{prop}_ridge_by_{group_col}.png", fig_dir)


def timecourse_by(records, prop, ylabel, datatype, fig_dir, *, time_col=None, roles=None,
                  series_col=None, colors=None) -> None:
    """Timecourse of per-sample means vs a unit-aware time axis, colored by `series_col`.
    Files: {datatype}_{prop}_timecourse[_{series_col}].png"""
    recs = _with_prop(records, prop)
    if not recs:
        return
    roles = roles or infer_roles(records)
    time_col = time_col if time_col is not None else _find_time_col(records, roles)
    if time_col is None:
        return
    trole = _role_for(records, time_col, roles)
    conv = trole.get("to_hours") or (lambda v: float(v))
    colors = colors or (color_map_for(records, series_col, roles) if series_col else {})

    series = {}
    for r in recs:
        try:
            x = conv(rget(r, time_col))
        except (TypeError, ValueError):
            continue
        key = rget(r, series_col) if series_col else "all"
        series.setdefault(key, []).append((x, r["props"][prop].mean()))

    fig, ax = plt.subplots(figsize=(8, 5))
    draw_timecourse(ax, series, colors, "Time (h)", ylabel)
    suffix = f"_{series_col}" if series_col else ""
    ax.set_title(f"{datatype} {prop} — timecourse" + (f" by {series_col}" if series_col else ""))
    _save(fig, f"{datatype}_{prop}_timecourse{suffix}.png", fig_dir)


def scatter_by(records, prop_x, prop_y, xlabel, ylabel, datatype, fig_dir, *, group_col=None,
               color_col=None, roles=None, ncols=4) -> None:
    """Per-cell scatter of prop_y vs prop_x with marginal histograms — a grid of per-sample panels.
    One figure per value of `group_col` (or a single figure if group_col is None). Panels colored
    by `color_col` (defaults to group_col). Pass records from load_ifxm_paired so x and y are
    row-aligned. Files: {datatype}_{prop_y}_vs_{prop_x}[_{group_col}={slug(value)}].png"""
    roles = roles or infer_roles(records)
    color_col = color_col or group_col
    colors = color_map_for(records, color_col, roles) if color_col else {}

    def _xy(r):
        x = np.asarray(r["props"].get(prop_x, []), float)
        y = np.asarray(r["props"].get(prop_y, []), float)
        if x.size == 0 or y.size == 0:
            return None
        if x.size != y.size:
            print(f"  WARNING: {r['sample']} {prop_x}/{prop_y} lengths differ "
                  f"({x.size} vs {y.size}) — skipping (use load_ifxm_paired?)")
            return None
        keep = np.isfinite(x) & np.isfinite(y)
        return (x[keep], y[keep]) if keep.sum() else None

    usable = [(r, xy) for r in records if (xy := _xy(r)) is not None]
    if not usable:
        return
    time_col = _find_time_col(records, roles)

    if group_col is None:
        groups = [(None, usable)]
    else:
        groups = [(v, [ru for ru in usable if rget(ru[0], group_col) == v])
                  for v in group_order([r for r, _ in usable], group_col, roles)]

    xy_by_sample = {r["sample"]: xy for r, xy in usable}
    for gval, sub in groups:
        if not sub:
            continue
        order_by = [c for c in (time_col, "rep") if c]
        ordered_recs = sort_records([r for r, _ in sub], order_by, roles) if order_by \
            else [r for r, _ in sub]
        sub = [(r, xy_by_sample[r["sample"]]) for r in ordered_recs]

        n = len(sub)
        nc = min(n, ncols)
        nrows = int(np.ceil(n / nc))
        fig = plt.figure(figsize=(nc * 3.5, nrows * 3.5))
        outer = gridspec.GridSpec(nrows, nc, figure=fig, hspace=0.55, wspace=0.45)
        for idx, (r, (x, y)) in enumerate(sub):
            ri, ci = divmod(idx, nc)
            c = colors.get(rget(r, color_col), FALLBACK_COLOR) if color_col else FALLBACK_COLOR
            draw_scatter_marginal(fig, outer[ri, ci], x, y, c, xlabel, ylabel,
                                  _detail_label(r, roles, time_col))
        for idx in range(n, nrows * nc):
            ri, ci = divmod(idx, nc)
            fig.add_subplot(outer[ri, ci]).axis("off")
        gsuffix = f"_{group_col}={_slug(gval)}" if group_col is not None else ""
        gtitle = f" — {value_label(group_col, gval, roles)}" if group_col is not None else ""
        fig.suptitle(f"{datatype} {prop_y} vs {prop_x}{gtitle}", fontsize=11)
        _save(fig, f"{datatype}_{prop_y}_vs_{prop_x}{gsuffix}.png", fig_dir)


def facet(records, prop, ylabel, datatype, fig_dir, *, facet_col, roles=None, inner="box",
          ncols=4, colors=None) -> None:
    """Compact grid: one panel per value of `facet_col` in a SINGLE figure (each panel pools that
    value's cells into one inner box/ridge). File: {datatype}_{prop}_facet_{facet_col}.png"""
    recs = _with_prop(records, prop)
    if not recs:
        return
    roles = roles or infer_roles(records)
    colors = colors or color_map_for(records, facet_col, roles)
    values = group_order(recs, facet_col, roles)
    n = len(values)
    if n == 0:
        return
    nc = min(n, ncols)
    nrows = int(np.ceil(n / nc))
    fig, axes = plt.subplots(nrows, nc, figsize=(nc * 3.2, nrows * 3.2), squeeze=False)
    for idx, v in enumerate(values):
        ax = axes[idx // nc][idx % nc]
        arr = np.concatenate([r["props"][prop] for r in recs if rget(r, facet_col) == v])
        c = colors.get(v, FALLBACK_COLOR)
        if inner == "ridge":
            draw_ridge(ax, [arr], [""], [c], ylabel)
        else:
            draw_boxes(ax, [arr], [""], [c], ylabel)
        ax.set_title(value_label(facet_col, v, roles), fontsize=9)
    for idx in range(n, nrows * nc):
        axes[idx // nc][idx % nc].axis("off")
    fig.suptitle(f"{datatype} {prop} — by {facet_col}", fontsize=11)
    _save(fig, f"{datatype}_{prop}_facet_{facet_col}.png", fig_dir)


def cross_groups(records, prop, ylabel, datatype, fig_dir, *, cols, roles=None,
                 kinds=("box", "ridge"), colors=None) -> None:
    """CROSSING on request: compare the cross-product of two (or more) columns. Adds a synthetic
    joined key (e.g. 'activated | DMEM') and runs compare_groups on it, per_sample.
    File: {datatype}_{prop}_{kind}_by_{colA}-x-{colB}.png"""
    cross_col = "-x-".join(cols)
    tagged = []
    for r in records:
        key = " | ".join(str(rget(r, c)) for c in cols)
        rr = dict(r)
        rr["meta"] = dict(r["meta"])
        rr["meta"][cross_col] = key
        tagged.append(rr)
    compare_groups(tagged, prop, ylabel, datatype, fig_dir, group_col=cross_col, roles=roles,
                   kinds=kinds, agg="per_sample", colors=colors)


# ---------------------------------------------------------------------------
# High-level: infer a plot plan, show it, execute it
# ---------------------------------------------------------------------------
def build_plan(records, datatype, *, roles=None, props=None, scatter_pairs=None,
               include_ordered=True) -> dict:
    """Infer a plot plan (list of PlotSpecs) from the column roles. For every boolean / categorical
    (and, if include_ordered, approved ordered) column: a plot_grouped + a compare_groups per prop.
    If a time column exists: a timecourse per prop, split by each grouping column. Scatters are
    added for scatter_pairs. `props` is filtered to those non-empty in at least one record, so a
    mass-only / volume-only experiment proposes no no-op plots for absent properties. Returns
    {roles, props, plots:[{fn,prop,kwargs,rationale}]}."""
    roles = roles or infer_roles(records)
    present = {k for r in records for k, v in r["props"].items() if len(v) > 0}
    props = [p for p in (props if props is not None else present) if p in present]
    time_col = _find_time_col(records, roles)

    group_cols = [c for c, i in roles.items()
                  if i["role"] in ("boolean", "categorical")
                  or (include_ordered and i["role"] == "ordered")]

    plots = []
    for prop in props:
        for gc in group_cols:
            plots.append({"fn": "plot_grouped", "prop": prop, "kwargs": {"group_col": gc},
                          "rationale": f"per-{gc} detail"})
            plots.append({"fn": "compare_groups", "prop": prop, "kwargs": {"group_col": gc},
                          "rationale": f"compare across {gc}"})
        if time_col:
            plots.append({"fn": "timecourse_by", "prop": prop,
                          "kwargs": {"time_col": time_col, "series_col": None},
                          "rationale": f"timecourse over {time_col}"})
            for gc in group_cols:
                plots.append({"fn": "timecourse_by", "prop": prop,
                              "kwargs": {"time_col": time_col, "series_col": gc},
                              "rationale": f"timecourse over {time_col}, split by {gc}"})
    for (px, py, xl, yl) in (scatter_pairs or []):
        if px not in present or py not in present:
            continue  # e.g. a mass-only experiment has no density/volume to scatter against
        gc = group_cols[0] if group_cols else None
        plots.append({"fn": "scatter_by", "prop": f"{py}_vs_{px}",
                      "kwargs": {"prop_x": px, "prop_y": py, "xlabel": xl, "ylabel": yl,
                                 "group_col": gc},
                      "rationale": "per-cell scatter" + (f" per {gc}" if gc else "")})
    return {"roles": roles, "props": props, "plots": plots}


def render_plan(plan) -> str:
    """Human-readable summary of inferred roles + proposed plots, with any confirm-me flags —
    the text to show the user for approval before generating/executing the driver."""
    roles = plan["roles"]
    lines = ["Inferred metadata roles:"]
    for col, i in roles.items():
        flag = "  [CONFIRM]" if i["confirm"] else ""
        unit = f" [{i['unit']}]" if i.get("unit") else ""
        lines.append(f"  - {col}: {i['role']}{unit} — {i['reason']}{flag}")
    confirms = [c for c, i in roles.items() if i["confirm"]]
    if confirms:
        lines.append("")
        lines.append("Needs your confirmation: " + ", ".join(confirms)
                     + "  (pass overrides={col: 'ordered'|'categorical'|...} to pin)")
    lines.append("")
    lines.append(f"Proposed plots ({len(plan['plots'])}):")
    for p in plan["plots"]:
        lines.append(f"  - {p['fn']}({p['prop']}) — {p['rationale']}")
    return "\n".join(lines)


_COMBINATORS = None


def autoplot(records, plan, datatype, fig_dir, prop_labels=None, paired_records=None) -> None:
    """Execute a plan's PlotSpecs. Distribution/timecourse plots run on `records`; scatter_by specs
    run on `paired_records` (row-aligned, from load_ifxm_paired) — pass it whenever the plan has
    scatter pairs, or those specs are skipped. `prop_labels`: {prop: axis_label} for y-axis labels."""
    global _COMBINATORS
    if _COMBINATORS is None:
        _COMBINATORS = {"plot_grouped": plot_grouped, "compare_groups": compare_groups,
                        "timecourse_by": timecourse_by, "scatter_by": scatter_by,
                        "facet": facet}
    labels = prop_labels or {}
    roles = plan["roles"]
    warned = False
    for spec in plan["plots"]:
        fn = _COMBINATORS[spec["fn"]]
        prop = spec["prop"]
        kw = dict(spec["kwargs"], roles=roles)
        if spec["fn"] == "scatter_by":
            if paired_records is None:
                if not warned:
                    print("  (skipping scatter specs — pass paired_records=load_ifxm_paired(...))")
                    warned = True
                continue
            fn(paired_records, datatype=datatype, fig_dir=fig_dir, **kw)
        else:
            fn(records, prop, labels.get(prop, prop), datatype, fig_dir, **kw)


# ---------------------------------------------------------------------------
# PowerPoint export
# ---------------------------------------------------------------------------
def save_pptx(fig_dir, out_path) -> None:
    """Compile every PNG in fig_dir into a 16:9 deck, one image per slide (centered, aspect-fit)."""
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.util import Inches, Emu

    fig_dir, out_path = Path(fig_dir), Path(out_path)
    prs = Presentation()
    slide_w, slide_h = Inches(13.33), Inches(7.5)
    prs.slide_width, prs.slide_height = slide_w, slide_h
    blank_layout = prs.slide_layouts[6]

    pngs = sorted(fig_dir.glob("*.png"))
    for png in pngs:
        img_w, img_h = PILImage.open(png).size
        if img_w / img_h > slide_w / slide_h:
            w, h = slide_w, Emu(int(slide_w * img_h / img_w))
        else:
            h, w = slide_h, Emu(int(slide_h * img_w / img_h))
        left = Emu(int((slide_w - w) // 2))
        top  = Emu(int((slide_h - h) // 2))
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(png), left=left, top=top, width=w, height=h)

    prs.save(str(out_path))
    print(f"  saved {out_path.name} ({len(pngs)} slides)")
